"""Bytes → text extractors for chat-uploaded files."""

from __future__ import annotations

from pathlib import Path


def _extract_pdf_text(raw: bytes) -> str:
    """Extract text from PDF bytes. Returns empty string on failure."""
    try:
        import io

        # Try PyMuPDF first (fast, no external deps)
        try:
            import fitz

            doc = fitz.open(stream=raw, filetype="pdf")
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
            return text.strip()
        except ImportError:
            pass
        # Fallback: pdfplumber
        try:
            import pdfplumber

            with pdfplumber.open(io.BytesIO(raw)) as pdf:
                return "\n".join(page.extract_text() or "" for page in pdf.pages).strip()
        except ImportError:
            pass
        return "[PDF text extraction not available — install PyMuPDF or pdfplumber]"
    except Exception as e:
        return f"[PDF extraction failed: {e}]"


def _extract_xlsx_text(raw: bytes) -> str:
    """Extract text from XLSX bytes as datatable JSON."""
    try:
        import io

        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            columns = [str(c) if c is not None else "" for c in rows[0]]
            data = [[str(c) if c is not None else "" for c in row] for row in rows[1:201]]
            parts.append(
                f"Sheet '{ws.title}': {len(rows) - 1} rows, {len(columns)} columns\n"
                + "\n".join(
                    f"{columns[i]}: {', '.join(r[i] for r in data[:5])}"
                    for i in range(min(len(columns), 10))
                )
            )
        wb.close()
        return "\n\n".join(parts) if parts else "[Empty spreadsheet]"
    except ImportError:
        return "[XLSX reading requires openpyxl — not installed]"
    except Exception as e:
        return f"[XLSX extraction failed: {e}]"


def _extract_docx_text(raw: bytes) -> str:
    """Extract text from DOCX bytes."""
    try:
        import io

        from docx import Document

        doc = Document(io.BytesIO(raw))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            headers = [c.text for c in table.rows[0].cells]
            parts.append(" | ".join(headers))
            for row in table.rows[1:]:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts) if parts else "[Empty document]"
    except ImportError:
        return "[DOCX reading requires python-docx — not installed]"
    except Exception as e:
        return f"[DOCX extraction failed: {e}]"


def _extract_pptx_text(raw: bytes) -> str:
    """Extract text from PPTX bytes."""
    try:
        import io

        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(cell.text for cell in row.cells))
            if texts:
                parts.append(f"Slide {i}:\n" + "\n".join(texts))
        return "\n\n".join(parts) if parts else "[Empty presentation]"
    except ImportError:
        return "[PPTX reading requires python-pptx — not installed]"
    except Exception as e:
        return f"[PPTX extraction failed: {e}]"


def _read_file_content(raw: bytes, mime: str, name: str, size: int) -> str:
    """Extract readable text from file bytes based on MIME type."""
    suffix = Path(name).suffix.lower()
    if suffix == ".xlsx" or suffix == ".xls" or "spreadsheet" in mime:
        return _extract_xlsx_text(raw) or f"[Empty spreadsheet, {size} bytes]"
    elif suffix == ".docx" or "wordprocessing" in mime:
        return _extract_docx_text(raw) or f"[Empty document, {size} bytes]"
    elif suffix == ".pptx" or "presentation" in mime:
        return _extract_pptx_text(raw) or f"[Empty presentation, {size} bytes]"
    elif any(
        t in mime
        for t in ["text/", "json", "xml", "csv", "html", "javascript", "yaml", "yml", "markdown"]
    ):
        return raw.decode("utf-8", errors="replace")
    elif "pdf" in mime:
        text = _extract_pdf_text(raw)
        if text:
            return text
        # No text layer extractable. This is the normal case for
        # scanned-image PDFs (CamScanner output, photographed
        return (
            f"[Scanned-image PDF, {size} bytes — no text layer to extract. "
            f"Use an OCR/vision tool (e.g. process_invoice for invoices) "
            f"to read the page contents from the image.]"
        )
    else:
        return f"[Binary file, {size} bytes]"
